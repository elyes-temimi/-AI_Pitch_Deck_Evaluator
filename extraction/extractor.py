"""
extraction/extractor.py
Normalizes any supported pitch deck input into a single canonical format:
    List[{"slide": int, "title": str, "content": str}]

Supported inputs:
- .json  -> structured deck, e.g. [{"slide":1,"title":"Problem","content":"..."}]
            (the challenge explicitly calls this out as acceptable:
            "slide-based text or structured format")
- .pptx  -> real PowerPoint file (title = title placeholder if present,
            content = remaining text, notes, tables)
- .pdf   -> one "slide" per page; title = first short line if plausible
- .txt/.md -> plain text, slides separated by a line containing only "---"

Keeping everything downstream in this one canonical shape means every agent
and the UI never needs to know which input format was used.
"""

from pathlib import Path
import json


def _guess_title(text: str, fallback: str):
    lines = [l for l in text.splitlines() if l.strip()]
    if not lines:
        return fallback, ""
    first = lines[0].strip()
    if len(first) <= 60:
        rest = "\n".join(lines[1:]).strip()
        return first, rest
    return fallback, text.strip()


def extract_json(path: Path):
    raw = json.loads(path.read_text(encoding="utf-8"))
    slides = []
    for i, item in enumerate(raw, start=1):
        slides.append({
            "slide": item.get("slide", i),
            "title": item.get("title", f"Slide {i}"),
            "content": item.get("content", ""),
        })
    return slides


def extract_pptx(path: Path):
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        body_chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                is_title_placeholder = (
                    shape.placeholder_format is not None
                    and shape.placeholder_format.type is not None
                    and "TITLE" in str(shape.placeholder_format.type)
                )
                if title is None and is_title_placeholder:
                    title = text
                else:
                    body_chunks.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells)
                    if row_text.strip(" |"):
                        body_chunks.append(row_text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                body_chunks.append(f"[Speaker notes: {note}]")
        content = "\n".join(body_chunks).strip()
        if title is None:
            title, content = _guess_title(content, f"Slide {i}")
        slides.append({"slide": i, "title": title, "content": content})
    return slides


def extract_pdf(path: Path):
    import pdfplumber

    slides = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            title, content = _guess_title(text, f"Slide {i}")
            slides.append({"slide": i, "title": title, "content": content})
    return slides


# --- Noise detection for low-quality / template / reproduction PDFs ---
#
# Real-world test (Airbnb's pitch deck via a third-party "PitchDeckCoach"
# reproduction PDF) showed extraction can pull in:
#   1. Literal Lorem Ipsum placeholder text left in an unfinished template slide
#   2. A repeated boilerplate footer line on nearly every slide
#   3. Trailing slides that are the TEMPLATE PROVIDER'S promotional content,
#      not part of the actual pitch deck at all
# Feeding any of this to the LLM as if it were the founder's actual content
# produces nonsensical clarity/narrative judgments (e.g. flagging Lorem Ipsum
# as "ambiguous phrasing"). This is caught and stripped/flagged here, BEFORE
# it reaches any agent — garbage-in/garbage-out is an extraction problem, not
# a prompting problem, and should be fixed at the extraction layer.


def extract_txt(path: Path):
    raw = path.read_text(encoding="utf-8", errors="ignore")
    parts = [s.strip() for s in raw.split("\n---\n") if s.strip()]
    slides = []
    for i, part in enumerate(parts, start=1):
        title, content = _guess_title(part, f"Slide {i}")
        slides.append({"slide": i, "title": title, "content": content})
    return slides


import re


def _strip_leading_trailing_number(title: str, slide_num: int):
    """Strips a slide-number artifact like 'Welcome 1' -> 'Welcome' when the
    PDF/PPTX extractor concatenated the page number onto the title text."""
    pattern = rf"^(.*?)\s*{slide_num}$"
    m = re.match(pattern, title.strip())
    if m and m.group(1).strip():
        return m.group(1).strip()
    return title


def _looks_like_placeholder_filler(text: str) -> bool:
    """Detects literal Lorem Ipsum placeholder text left in a template the
    user is evaluating. This is a DATA-QUALITY issue with the source file,
    not an extraction bug — flagged explicitly rather than silently scored
    as if it were real founder writing."""
    return bool(re.search(r"lorem ipsum", text, re.IGNORECASE))


_PROMO_SLIDE_PATTERNS = [
    r"slideshare\.net", r"did you enjoy this deck", r"you'?ll love our",
    r"read my blog post", r"ultimate pitch deck template",
]

# Generic footer/watermark pattern: "Template by <Anything>.com" can appear
# INLINE within real slide content (not always on its own line), so it must
# be stripped by regex substitution, not just by exact-line frequency match.
_INLINE_FOOTER_PATTERN = re.compile(r"\btemplate\s+by\s+[\w\.\-]+\.com\b", re.IGNORECASE)


def _strip_inline_footer(text: str) -> str:
    return _INLINE_FOOTER_PATTERN.sub("", text).strip()


def _looks_like_promo_slide(title: str, content: str) -> bool:
    """Detects trailing slides that are template-vendor advertising
    (e.g. 'check out our other templates') rather than actual pitch
    content. These should not count toward narrative/clarity scoring."""
    blob = f"{title}\n{content}".lower()
    return any(re.search(p, blob) for p in _PROMO_SLIDE_PATTERNS)


def clean_slides(slides):
    """
    Post-extraction cleanup for PDF/PPTX/TXT sources (not applied to .json
    input, which is already structured by the user). Fixes three concrete
    issues observed when extracting real-world template-sourced decks:

    1. Slide-number artifacts concatenated onto titles ("Welcome 1" -> "Welcome").
    2. Repeated boilerplate footer/watermark lines (e.g. "Template by X.com")
       that appear on most slides — auto-detected by frequency, not hardcoded
       to one vendor, so this generalizes to other watermarked templates.
    3. Trailing non-content slides (template-vendor ads, "click here" links)
       that aren't part of the actual pitch and would distort narrative scoring.

    Also flags (but does not strip) literal placeholder filler text
    (Lorem Ipsum), since that's a data-quality problem with the source file
    itself, not something safe to silently paper over.

    Returns (cleaned_slides, warnings).
    """
    warnings = []
    if not slides:
        return slides, warnings

    # --- detect repeated boilerplate lines by frequency ---
    line_counts = {}
    for s in slides:
        seen_this_slide = set()
        for line in s["content"].splitlines():
            line_norm = line.strip()
            if line_norm and line_norm not in seen_this_slide:
                line_counts[line_norm] = line_counts.get(line_norm, 0) + 1
                seen_this_slide.add(line_norm)
    boilerplate_lines = {
        line for line, count in line_counts.items()
        if len(slides) >= 3 and count / len(slides) >= 0.4
    }

    cleaned = []
    dropped_promo = []
    placeholder_flagged = []
    inline_footer_hits = 0

    for s in slides:
        title = _strip_leading_trailing_number(s["title"], s["slide"])
        content_lines = [
            line for line in s["content"].splitlines()
            if line.strip() not in boilerplate_lines
        ]
        content = "\n".join(content_lines).strip()

        before = content
        content = _strip_inline_footer(content)
        if content != before:
            inline_footer_hits += 1

        if _looks_like_promo_slide(title, content):
            dropped_promo.append(f"Slide {s['slide']} ({title})")
            continue  # excluded entirely — not pitch content

        if _looks_like_placeholder_filler(content):
            placeholder_flagged.append(f"Slide {s['slide']} ({title})")

        cleaned.append({"slide": s["slide"], "title": title, "content": content})

    if boilerplate_lines or inline_footer_hits:
        warnings.append(
            f"Stripped {len(boilerplate_lines)} repeated boilerplate/watermark "
            f"line(s) and {inline_footer_hits} inline footer mention(s) "
            f"(e.g. 'Template by X.com' appearing inside slide text) before evaluation."
        )
    if dropped_promo:
        warnings.append(
            f"Excluded {len(dropped_promo)} trailing slide(s) that look like "
            f"template-vendor advertising rather than pitch content: {', '.join(dropped_promo)}"
        )
    if placeholder_flagged:
        warnings.append(
            f"Literal 'Lorem ipsum' placeholder filler text detected on: "
            f"{', '.join(placeholder_flagged)}. This is a data-quality issue with "
            f"the source file itself (unfilled template content) — the evaluator "
            f"will score this content as written, but it should not be read as a "
            f"genuine flaw in the original pitch."
        )

    return cleaned, warnings


def extract_deck(path_or_bytes, filename: str = None):
    """Accepts a filesystem path, or (bytes, filename) from a file uploader.
    Returns (slides, warnings)."""
    if isinstance(path_or_bytes, (str, Path)):
        p = Path(path_or_bytes)
        suffix = p.suffix.lower()
        source_path = p
    else:
        import tempfile
        suffix = Path(filename).suffix.lower()
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        tmp.write(path_or_bytes)
        tmp.close()
        source_path = Path(tmp.name)

    if not source_path.exists():
        raise FileNotFoundError(f"File not found: {source_path}")

    if suffix == ".json":
        slides = extract_json(source_path)
    elif suffix == ".pptx":
        slides = extract_pptx(source_path)
    elif suffix == ".pdf":
        slides = extract_pdf(source_path)
    elif suffix in (".txt", ".md"):
        slides = extract_txt(source_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix} (use .json, .pptx, .pdf, or .txt)")

    slides, cleaning_warnings = clean_slides(slides)

    warnings = list(cleaning_warnings)
    thin = [s for s in slides if len(s["content"]) < 15]
    if slides and len(thin) / len(slides) > 0.3:
        warnings.append(
            f"{len(thin)}/{len(slides)} slides have almost no extractable text "
            "(likely image-only/scanned content). This MVP does not run OCR, "
            "so evaluation quality on those slides will be limited."
        )
    if not slides:
        warnings.append("No slides could be read from this file.")

    return slides, warnings


def deck_to_prompt_text(slides) -> str:
    lines = []
    for s in slides:
        content = s["content"] if s["content"].strip() else "(no extractable text)"
        lines.append(f"--- Slide {s['slide']}: {s['title']} ---\n{content}")
    return "\n\n".join(lines)


def get_section(slides, *keywords):
    """Finds the first slide whose title contains any of the given keywords
    (case-insensitive). Used by the calibration agent to pull out, e.g., the
    'Problem' slide's text to use as a retrieval query. Returns '' if none
    found, so callers can gracefully fall back to the full deck text."""
    for s in slides:
        title_lower = s["title"].lower()
        if any(k.lower() in title_lower for k in keywords):
            return s["content"]
    return ""